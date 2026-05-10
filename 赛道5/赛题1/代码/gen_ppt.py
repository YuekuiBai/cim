#!/usr/bin/env python3
"""
自动生成CIM比赛PPT - 使用实验图表
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

RESULTS_DIR = '/mnt/storage2/zyc/CIM比赛/赛道5/赛题1/results'
OUTPUT_PATH = '/mnt/storage2/zyc/CIM比赛/赛道5/赛题1/文档/介绍PPT.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 112, 192)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(89, 89, 89)

def add_cover_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 248, 255)
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    team_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.6))
    tf = team_box.text_frame
    p = tf.paragraphs[0]
    p.text = "团队：HDUer"
    p.font.size = Pt(24)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_title_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_text_slide(prs, title, lines):
    add_title_slide(prs, title)
    slide = prs.slides[-1]

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.font.color.rgb = GRAY
        p.space_after = Pt(14)

def add_image_slide(prs, title, image_path, description=""):
    add_title_slide(prs, title)
    slide = prs.slides[-1]

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(11))

    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(12), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

def add_table_slide(prs, title, headers, rows, col_widths=None):
    add_title_slide(prs, title)
    slide = prs.slides[-1]

    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)).table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            if col_idx > 0:
                p.alignment = PP_ALIGN.CENTER

# ==================== 开始生成PPT ====================

# 第1页：封面
add_cover_slide(prs,
    "存算一体芯片中非线性误差对推理精度的影响研究",
    "赛道五：存算算法")

# 第2页：目录
add_text_slide(prs, "目录", [
    "1. 研究背景",
    "2. 敏感性分析",
    "3. 鲁棒性增强方法",
    "4. 实验结果",
    "5. 拓展研究",
    "6. 总结与展望"
])

# 第3页：研究背景 - 章节页
add_section_slide(prs, "1. 研究背景")

# 第4页：研究背景 - 内容
add_text_slide(prs, "1. 研究背景", [
    "存算一体(CIM)架构通过消除数据搬运提升能效",
    "",
    "• 传统冯诺依曼架构面临存储墙瓶颈",
    "• CIM将计算嵌入存储单元，减少数据移动",
    "• 模拟器件实现的CIM易受非线性误差影响",
    "",
    "主要挑战：",
    "• 输入信号幅度变化导致非线性特性",
    "• 温度漂移和器件工艺偏差"
])

# 第5页：研究问题
add_text_slide(prs, "1.1 研究问题", [
    "核心问题：",
    "如何通过算法增强神经网络在实际存算芯片部署中",
    "对非线性误差的容忍能力与鲁棒性？",
    "",
    "研究任务：",
    "• 任务一：分析不同非线性强度下推理精度的衰减特性",
    "• 任务二：设计与实现非线性感知训练方法",
    "• 任务三：评估多种鲁棒性增强方法的效果"
])

# 第6页：敏感性分析 - 章节页
add_section_slide(prs, "2. 敏感性分析")

# 第7页：敏感性分析 - 图表
add_image_slide(prs, "2.1 非线性强度对推理精度的影响",
    f"{RESULTS_DIR}/task1_sensitivity/accuracy_vs_alpha.png",
    "ResNet18在CIFAR-10上，不同α值对应的推理精度")

# 第8页：敏感性分析 - 详细数据
add_table_slide(prs, "2.2 敏感性分析详细结果",
    ["α值", "推理精度", "精度衰减", "敏感等级"],
    [
        ["0.0 (Clean)", "81.95%", "基准", "无"],
        ["0.1", "80.99%", "-1.0%", "低"],
        ["0.2", "77.63%", "-4.3%", "开始下降"],
        ["0.3", "70.95%", "-11.0%", "显著损失"],
        ["0.4", "60.40%", "-21.5%", "严重下降"],
        ["0.5", "48.49%", "-33.5%", "接近一半损失"]
    ])

# 第9页：误差累积分析
add_image_slide(prs, "2.3 误差累积分析",
    f"{RESULTS_DIR}/task1_sensitivity/error_accumulation.png",
    "各层非线性误差的累积效应")

# 第10页：层分布偏移
add_image_slide(prs, "2.4 层分布偏移分析",
    f"{RESULTS_DIR}/task1_sensitivity/layer_distribution_shift.png",
    "不同层激活分布随非线性强度增加的变化")

# 第11页：鲁棒性方法 - 章节页
add_section_slide(prs, "3. 鲁棒性增强方法")

# 第12页：方法概述
add_text_slide(prs, "3.1 鲁棒性增强方法概述", [
    "本研究评估了多种鲁棒性增强方法：",
    "",
    "1. 预失真补偿 (Pre-distortion)",
    "   - 在前向传播前对输入进行非线性预补偿",
    "",
    "2. 校准层 (Calibration Layer)",
    "   - 添加可学习的校准层修正非线性输出",
    "",
    "3. 非线性感知训练 (NAT)",
    "   - 在训练时注入非线性噪声",
    "",
    "4. OVF训练 (Oriented Variational Forward)",
    "   - 基于负反馈理论，定向噪声注入"
])

# 第13页：NAT方法详解
add_text_slide(prs, "3.2 非线性感知训练 (NAT)", [
    "核心思想：",
    "在训练过程中注入非线性误差，使网络学习适应非线性环境",
    "",
    "实现方式：",
    "• 前向传播：x' = x + α(x³ - x)",
    "• 损失计算：在非线性扰动后计算损失",
    "• 反向传播：更新权重以抵抗非线性影响",
    "",
    "NAT+混合Alpha优化：",
    "• 每次迭代随机采样α∈[0, 0.5]",
    "• 扩大训练分布覆盖范围"
])

# 第14页：OVF和SAM方法
add_text_slide(prs, "3.3 OVF与SAM训练方法", [
    "OVF (Oriented Variational Forward) 训练",
    "• 基于负反馈理论，定向注入变异噪声",
    "• 在CODES+ISSS 2024发表",
    "",
    "SAM (Sharpness-Aware Minimization) 训练",
    "• 最小化损失面的尖锐度",
    "• 提升模型泛化能力和鲁棒性",
    "• 在ICLR 2021发表"
])

# 第15页：实验结果 - 章节页
add_section_slide(prs, "4. 实验结果")

# 第16页：方法对比总览图
add_image_slide(prs, "4.1 鲁棒性方法综合对比",
    f"{RESULTS_DIR}/figures/summary_comparison.png",
    "各方法在不同非线性强度下的精度表现")

# 第17页：方法对比表格
add_table_slide(prs, "4.2 方法综合对比表",
    ["方法", "α=0.0", "α=0.2", "α=0.5", "平均", "波动"],
    [
        ["基线", "81.95%", "77.63%", "48.49%", "70.07%", "33.46%"],
        ["预失真补偿", "81.95%", "79.71%", "60.74%", "74.68%", "21.21%"],
        ["NAT+混合Alpha", "79.03%", "83.74%", "78.56%", "81.71%", "5.53%"],
        ["OVF训练", "77.40%", "83.47%", "78.11%", "81.06%", "6.24%"],
        ["SAM训练", "76.45%", "82.29%", "77.44%", "80.06%", "6.30%"]
    ])

# 第18页：Scratch vs Finetune
add_image_slide(prs, "4.3 微调 vs 从头训练",
    f"{RESULTS_DIR}/task2_training/scratch_rigorous/comparison_chart.png",
    "微调策略显著优于从头训练 (p < 0.0001)")

# 第19页：量化与非线性对比
add_image_slide(prs, "4.4 量化 vs 非线性误差",
    f"{RESULTS_DIR}/figures/quantization_vs_nonlinear.png",
    "量化误差与非线性误差的独立处理")

# 第20页：拓展研究 - 章节页
add_section_slide(prs, "5. 拓展研究")

# 第21页：网络结构影响
add_text_slide(prs, "5.1 网络结构影响", [
    "ResNet34最鲁棒（衰减仅3.79%）",
    "",
    "实验结果：",
    "• ResNet34: 衰减3.79% - 最优",
    "• ResNet18: 衰减4.48% - 中等",
    "• MobileNetV2: 衰减5.37% - 最敏感",
    "",
    "结论：",
    "• 残差连接优于深度可分离卷积",
    "• 适度增加网络深度可提升鲁棒性"
])

# 第22页：高斯噪声等效
add_text_slide(prs, "5.2 高斯噪声等效分析", [
    "发现：非线性误差与高斯噪声有等效关系",
    "",
    "等效关系：",
    "• σ=0.1 ≈ α=0.2",
    "• σ=0.15 ≈ α=0.4",
    "",
    "应用价值：",
    "• 可用高斯噪声简化非线性实验",
    "• 便于理论分析和建模"
])

# 第23页：主要贡献
add_text_slide(prs, "5.3 主要贡献", [
    "1. 系统性分析",
    "   对不同非线性强度下神经网络推理精度的衰减特性进行了全面分析",
    "",
    "2. 方法对比",
    "   对比研究了多种鲁棒性增强方法",
    "",
    "3. 创新方法",
    "   提出了NAT+混合Alpha优化方法",
    "",
    "4. 严谨实验",
    "   通过统计学方法验证微调策略相对于从头训练的优势"
])

# 第24页：总结与展望 - 章节页
add_section_slide(prs, "6. 总结与展望")

# 第25页：研究成果
add_text_slide(prs, "6.1 研究成果", [
    "研究成果：",
    "",
    "• 系统分析了CIM中非线性误差对网络精度的影响规律",
    "• 验证了NAT系列方法的有效性",
    "• 提出了NAT+混合Alpha优化，取得了最佳的鲁棒性-精度平衡",
    "• NAT+混合Alpha平均精度81.71%，波动仅±5.53%"
])

# 第26页：未来方向
add_text_slide(prs, "6.2 未来方向", [
    "未来研究方向：",
    "",
    "• 在更大规模数据集(ImageNet)上验证",
    "• 结合量化误差与非线性误差的联合优化",
    "• 探索更多先进的训练策略(如SAM变体)",
    "• 实际硬件平台部署验证"
])

# 第27页：团队介绍
add_text_slide(prs, "7. 团队介绍", [
    "团队名称：HDUer",
    "",
    "工作说明：",
    "独立完成本项目的全部研究工作",
    "",
    "个人职责：",
    "• 项目统筹、算法设计",
    "• 实验实现、报告撰写",
    "",
    "致谢：",
    "感谢赛道五主办方提供的宝贵竞赛机会"
])

# 第28页：结束页
add_cover_slide(prs, "感谢聆听", "欢迎批评指正")

# 保存
prs.save(OUTPUT_PATH)
print(f"PPT已更新: {OUTPUT_PATH}")