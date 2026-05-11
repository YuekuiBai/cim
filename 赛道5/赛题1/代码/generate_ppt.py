#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, content_blocks, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    y_pos = 1.2
    for block in content_blocks:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(11.5), Inches(0.4))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = block
        p.font.size = Pt(18)
        y_pos += 0.4

    if image_path and os.path.exists(image_path):
        img = slide.shapes.add_picture(image_path, Inches(1), Inches(4), width=Inches(11))
    return slide

def create_problem1_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    base_path = "/mnt/storage2/zyc/CIM比赛/赛道5/赛题1"
    figures_path = os.path.join(base_path, "results", "figures")
    task1_path = os.path.join(base_path, "results", "task1_sensitivity")

    add_title_slide(prs, "存算一体芯片中非线性误差对推理精度的影响研究", "赛道五：存算算法 - 赛题一\nHDUer团队")

    add_content_slide(prs, "目录", [
        "1. 研究背景与意义",
        "2. 非线性误差建模",
        "3. 三大任务设计",
        "4. 实验结果与分析",
        "5. 核心创新点",
        "6. 总结与展望"
    ])

    add_content_slide(prs, "研究背景与意义", [
        "存算一体技术：传统冯诺依曼架构面临存储墙瓶颈，存算一体将计算嵌入存储器件内部",
        "非线性误差问题：模拟器件易受环境因素干扰，导致乘积累加操作呈现非线性特性",
        "研究目标：如何通过算法增强神经网络对非线性误差的容忍能力与鲁棒性？"
    ])

    add_content_slide(prs, "非线性误差建模", [
        "数学模型：x' = α · x³ + (1 - α) · x",
        "参数说明：x为理想输入值，x'为非线性失真后的实际信号，α为非线性强度",
        "特性：α=0为理想线性关系；α>0时正输入缩小，负输入放大"
    ])

    add_content_slide(prs, "任务一：敏感性分析", [
        "实验设置：ResNet18（CIFAR-10），基准精度81.95%（α=0），测试α范围0.0~0.5",
        "关键发现：α<0.2影响可接受；α>0.3精度急剧下降；α=0.5精度损失超30%"
    ], os.path.join(task1_path, "accuracy_vs_alpha.png"))

    add_content_slide(prs, "敏感性分析详图", [
        "左图：误差累积随层深变化",
        "右图：层分布偏移分析"
    ], os.path.join(task1_path, "error_accumulation.png"))

    add_content_slide(prs, "任务二：非线性感知训练（NAT）", [
        "基本原理：训练阶段注入非线性误差，使模型学习适应非线性环境",
        "感知训练结果：α=0.2时，Clean精度84.34%，指定噪声下精度87.30%"
    ])

    add_content_slide(prs, "微调 vs 从头训练对比", [
        "实验结果：微调平均85.21%，从头训练81.32%",
        "差异：+3.85%~+4.15%，p-value<0.0001",
        "结论：微调策略优势显著"
    ])

    add_content_slide(prs, "任务三：鲁棒性增强方法", [
        "基线：81.95%(α=0) → 48.49%(α=0.5)，平均70.07%",
        "预失真补偿：平均74.68%",
        "校准层：平均76.73%",
        "NAT：平均78.25%，指定噪声下可达87.30%"
    ])

    add_content_slide(prs, "高级方法探索", [
        "NAT+混合Alpha：全范围均衡，平均81.71%，波动仅5.53%",
        "课程NAT：高α区域保护效果好，平均79.47%",
        "OVF训练：基于负反馈理论，平均81.06%，波动6.24%"
    ])

    add_content_slide(prs, "实验总览", [], os.path.join(figures_path, "experiment_overview.png"))

    add_content_slide(prs, "方法对比总结", [
        "噪声水平已知且固定 → NAT（对应α），效果>87%",
        "噪声水平未知/变化 → NAT+混合Alpha，全范围均衡",
        "噪声水平较高(α>0.3) → 课程NAT，高α保护好",
        "需要快速部署 → 预失真补偿，无需重训练"
    ])

    add_content_slide(prs, "核心创新点", [
        "创新一：NAT+混合Alpha优化方法，全范围均衡精度81.71%，波动最小5.53%",
        "创新二：课程NAT方法，借鉴课程学习思想，渐进式训练",
        "创新三：严谨实验验证，3×3独立实验设计，p-value<0.0001",
        "创新四：系统方法对比，覆盖10+种鲁棒性方法"
    ])

    add_content_slide(prs, "实验结论", [
        "敏感性分析：非线性误差影响呈非线性关系，α>0.3时精度急剧下降",
        "微调策略优势：显著优于从头训练（+4%），p-value<0.0001",
        "NAT+混合Alpha最优：全范围保护，平均81.71%，波动仅5.53%",
        "未来工作：更大规模网络验证，ImageNet测试，与CIM硬件协同设计"
    ])

    add_title_slide(prs, "感谢聆听", "赛道五：存算算法 - 赛题一\nHDUer团队\n欢迎提问与交流")

    output_path = os.path.join(base_path, "文档", "赛题一PPT.pptx")
    prs.save(output_path)
    print(f"赛题一PPT已生成: {output_path}")

if __name__ == "__main__":
    create_problem1_ppt()
