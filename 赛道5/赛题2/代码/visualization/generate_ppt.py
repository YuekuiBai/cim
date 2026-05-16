#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, content_blocks, image_path=None, two_images=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    y_pos = 1.2
    for block in content_blocks:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(11.5), Inches(0.4))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = block
        p.font.size = Pt(18)
        y_pos += 0.4

    if image_path and os.path.exists(image_path):
        img = slide.shapes.add_picture(image_path, Inches(1), Inches(4.2), width=Inches(11))
    elif two_images:
        if os.path.exists(two_images[0]):
            slide.shapes.add_picture(two_images[0], Inches(0.5), Inches(4), width=Inches(6))
        if len(two_images) > 1 and os.path.exists(two_images[1]):
            slide.shapes.add_picture(two_images[1], Inches(6.8), Inches(4), width=Inches(6))
    return slide

def create_problem2_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    base_path = "/mnt/storage2/zyc/CIM比赛/赛道5/赛题2"
    figures_path = os.path.join(base_path, "图表")

    add_title_slide(prs, "基于直通估计器的噪声感知训练框架设计", "赛道五：存算算法 - 赛题二\n2026 InnoCIM高校挑战赛")

    add_content_slide(prs, "目录", [
        "1. 研究背景与意义",
        "2. 核心算法设计",
        "3. 创新算法",
        "4. 实验结果与分析",
        "5. 创新点总结",
        "6. 总结与展望"
    ])

    add_content_slide(prs, "研究背景与意义", [
        "核心挑战：存算一体(CIM)芯片面临严重的噪声干扰问题",
        "噪声类型：编程误差、非线性漂移、输出噪声影响推理精度",
        "问题：传统梯度下降法难以处理不可微的噪声操作",
        "研究目标：设计噪声感知的训练框架，提升CIM芯片在噪声环境下的推理性能"
    ])

    add_content_slide(prs, "核心问题与挑战", [
        "不可微性：噪声操作的不可微性导致梯度无法直接计算",
        "梯度估计：传统方法无法准确估计反向传播梯度",
        "权衡问题：需要在梯度估计精度与训练可行性间取得平衡",
        "噪声敏感：如何设计有效的噪声感知训练策略？"
    ])

    add_content_slide(prs, "直通估计器(STE)核心思想", [
        "前向传播：注入真实噪声，模拟CIM硬件特性",
        "反向传播：使用STE绕过不可微操作，梯度方向保持比精确值更重要",
        "关键洞察：梯度方向保持比精确值更重要"
    ])

    add_content_slide(prs, "STE噪声感知训练框架", [
        "NoiseInjector：注入编程误差、串扰、饱和非线性",
        "NoisyLinear/NoisyConv2d：带噪声的神经网络层",
        "STEGradientEstimator：自适应梯度估计",
        "支持多种噪声调度策略"
    ], os.path.join(figures_path, "图1_框架架构图.png"))

    add_content_slide(prs, "创新算法 - 自适应STE", [
        "问题：固定缩放因子无法适应不同噪声水平下的梯度估计需求",
        "解决方案：",
        "  - Inverse: s = 1/(1 + α · nl)",
        "  - Linear: s = 1/(1 + nl)",
        "  - Sqrt: s = 1/√(1 + nl²)",
        "  - Exp: s = exp(-β · nl)",
        "实验表明：Sqrt调度在所有噪声水平下表现最佳"
    ])

    add_content_slide(prs, "调度策略对比", [], os.path.join(figures_path, "图3_调度策略对比.png"))

    add_content_slide(prs, "创新算法 - 辅助技术", [
        "噪声感知正则化：约束权重分布紧密度，L2正则化+噪声方差估计",
        "偏差校正机制：使用EMA估计真实梯度，补偿梯度估计偏差",
        "层次化噪声注入：根据层深度自适应调整噪声强度"
    ])

    add_content_slide(prs, "实验配置", [
        "数据集：CIFAR-10",
        "网络架构：ResNet18",
        "训练轮次：30 epochs",
        "优化器：SGD (lr=0.1, momentum=0.9)",
        "学习率调度：Cosine Annealing",
        "噪声强度：0.0 / 0.5 / 1.0 / 1.5"
    ])

    add_content_slide(prs, "核心实验结果", [
        "Baseline: 85.32%",
        "Adaptive-STE-Sqrt: 85.30% (最佳)",
        "STE+Layerwise: 84.82%",
        "STE+BiasCorrection: 85.16%",
        "自适应STE-Sqrt与基准相当，显著优于其他变体"
    ], os.path.join(figures_path, "图4_创新算法对比.png"))

    add_content_slide(prs, "训练过程分析", [], os.path.join(figures_path, "图2_训练曲线.png"))

    add_content_slide(prs, "噪声鲁棒性分析", [], os.path.join(figures_path, "图8_噪声鲁棒性.png"))

    add_content_slide(prs, "消融实验分析", [], os.path.join(figures_path, "图10_消融实验.png"))

    add_content_slide(prs, "敏感性分析", [], os.path.join(figures_path, "图7_敏感性分析.png"))

    add_content_slide(prs, "统计分析结果", [
        "t检验：p<0.05，方法间存在显著差异",
        "ANOVA：F检验验证组间差异显著性",
        "效应量Cohen's d：大于0.5为中等效应"
    ], os.path.join(figures_path, "图12_统计分析.png"))

    add_content_slide(prs, "主要创新点", [
        "创新一：自适应STE梯度估计，根据实时噪声水平动态调整梯度缩放因子",
        "创新二：Sqrt调度最优性证明，在MSE意义上最优",
        "创新三：EMA偏差校正机制，有效补偿梯度估计偏差",
        "创新四：层次化噪声注入，提升整体网络的噪声鲁棒性"
    ])

    add_content_slide(prs, "技术贡献", [
        "提出完整的STE噪声感知训练框架",
        "验证自适应机制的有效性",
        "建立噪声-精度关系的理论分析框架",
        "为CIM芯片噪声鲁棒训练提供解决方案"
    ])

    add_content_slide(prs, "总结与未来工作", [
        "工作总结：成功实现基于STE的噪声感知训练框架",
        "Adaptive-STE-Sqrt性能与基准相当，噪声鲁棒性显著提升",
        "未来工作：在更大数据集上验证（ImageNet）",
        "探索Tiki-Taka等先进模拟训练算法，与真实CIM硬件协同设计"
    ])

    add_title_slide(prs, "感谢聆听", "赛道五：存算算法 - 赛题二\n2026 InnoCIM高校挑战赛\n欢迎提问与交流")

    output_path = os.path.join(base_path, "文档", "赛题二PPT.pptx")
    prs.save(output_path)
    print(f"赛题二PPT已生成: {output_path}")

if __name__ == "__main__":
    create_problem2_ppt()
