#!/usr/bin/env python3
"""
生成赛题二答辩PPT（包含实际图表）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

FIGURES_DIR = "/mnt/storage2/zyc/CIM比赛/赛道5/赛题2/图表"

def add_title_slide(prs, title_text, subtitle_text=None):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title_text, bullet_points):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(24)
        p.space_after = Pt(16)

    return slide

def add_image_slide(prs, title_text, image_path, caption=None):
    """添加图片幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True

    if os.path.exists(image_path):
        img = slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(11))
    else:
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[图片: {os.path.basename(image_path)}]"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
        p.font.italic = True

    return slide

def add_text_slide(prs, title_text, content_text):
    """添加文本内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(22)
        p.space_after = Pt(8)

    return slide

def add_two_images_slide(prs, title_text, image_path1, image_path2, caption=None):
    """添加双图片幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.bold = True

    if os.path.exists(image_path1):
        slide.shapes.add_picture(image_path1, Inches(0.3), Inches(0.8), width=Inches(6.3))
    if os.path.exists(image_path2):
        slide.shapes.add_picture(image_path2, Inches(6.7), Inches(0.8), width=Inches(6.3))

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        p.font.italic = True

    return slide

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== 第1页：封面 ==========
    add_title_slide(prs,
        "基于直通估计器的噪声感知训练框架设计",
        "赛道五：存算算法 - 赛题二\n2026 InnoCIM高校挑战赛")

    # ========== 第2页：目录 ==========
    add_content_slide(prs, "汇报提纲", [
        "研究背景与意义",
        "核心算法设计",
        "实验结果与分析",
        "创新点总结",
        "总结与展望"
    ])

    # ========== 第3页：研究背景 ==========
    add_content_slide(prs, "研究背景与意义", [
        "存算一体(CIM)芯片面临噪声挑战",
        "编程误差、非线性噪声、输出误差影响推理精度",
        "传统梯度下降法难以处理不可微噪声操作",
        "需要设计噪声感知的训练框架"
    ])

    # ========== 第4页：框架架构图 ==========
    add_image_slide(prs, "STE噪声感知训练框架架构",
        f"{FIGURES_DIR}/图1_框架架构图.png",
        "图1：STE噪声感知训练框架整体架构")

    # ========== 第5页：问题分析 ==========
    add_content_slide(prs, "核心问题与挑战", [
        "噪声操作的不可微性导致梯度无法直接计算",
        "传统方法无法准确估计反向传播梯度",
        "需要在\"梯度估计精度\"与\"训练可行性\"间取得平衡",
        "如何设计有效的噪声感知训练策略？"
    ])

    # ========== 第6页：自适应调度对比 ==========
    add_image_slide(prs, "自适应STE调度策略对比",
        f"{FIGURES_DIR}/图3_调度策略对比.png",
        "图2：四种调度策略的缩放因子随噪声水平的变化曲线，Sqrt调度与理论最优最吻合")

    # ========== 第7页：创新算法（一） ==========
    add_content_slide(prs, "创新算法 - 自适应STE", [
        "问题：固定缩放因子无法适应不同噪声水平",
        "解决：设计自适应梯度缩放策略",
        "四种调度：Inverse / Linear / Sqrt / Exp",
        "Sqrt调度在实验中表现最佳"
    ])

    # ========== 第8页：创新算法（二） ==========
    add_content_slide(prs, "创新算法 - 噪声感知正则化", [
        "思想：约束权重分布紧密度，降低噪声敏感性",
        "方法：L2正则化 + 噪声方差估计",
        "偏差校正：EMA方法估计并补偿梯度偏差",
        "层次化注入：根据层深度自适应调整噪声强度"
    ])

    # ========== 第9页：实验设置 ==========
    add_text_slide(prs, "实验配置", """• 数据集：CIFAR-10图像分类
• 网络架构：ResNet18
• 训练配置：30 epochs, SGD, Cosine LR
• 噪声强度：0.0 / 0.5 / 1.0 / 1.5
• 对比方法：Baseline / STE-NAT / Adaptive-STE等""")

    # ========== 第10页：创新算法对比图 ==========
    add_image_slide(prs, "创新算法性能对比",
        f"{FIGURES_DIR}/图4_创新算法对比.png",
        "图3：各创新算法在CIFAR-10上的精度对比，Adaptive-STE-Sqrt达到最高85.30%")

    # ========== 第11页：训练曲线 ==========
    add_image_slide(prs, "训练过程曲线",
        f"{FIGURES_DIR}/图2_训练曲线.png",
        "图4：各方法的训练损失和精度随epoch的变化曲线")

    # ========== 第12页：消融实验 ==========
    add_image_slide(prs, "消融实验分析",
        f"{FIGURES_DIR}/图10_消融实验.png",
        "图5：各组件的独立贡献，Layerwise和Sqrt调度是关键改进")

    # ========== 第13页：噪声鲁棒性 ==========
    add_image_slide(prs, "噪声鲁棒性分析",
        f"{FIGURES_DIR}/图8_噪声鲁棒性.png",
        "图6：噪声训练与干净训练的对比，以及不同方法的抗噪能力")

    # ========== 第14页：敏感性分析 ==========
    add_image_slide(prs, "噪声敏感性分析",
        f"{FIGURES_DIR}/图7_敏感性分析.png",
        "图7：基准模型与STE-NAT在不同噪声水平下的性能变化")

    # ========== 第15页：统计显著性 ==========
    add_image_slide(prs, "统计显著性分析",
        f"{FIGURES_DIR}/图12_统计分析.png",
        "图8：箱线图显示方法稳定性，t检验验证显著性差异(p<0.05)")

    # ========== 第16页：创新点总结 ==========
    add_content_slide(prs, "主要创新点", [
        "自适应STE梯度估计：根据噪声水平动态调整缩放因子",
        "多调度策略验证：Sqrt调度效果最佳",
        "偏差校正机制：EMA方法补偿梯度估计偏差",
        "层次化噪声注入：根据层特性差异化处理"
    ])

    # ========== 第17页：技术贡献 ==========
    add_content_slide(prs, "技术贡献", [
        "提出完整的STE噪声感知训练框架",
        "验证自适应机制的有效性",
        "建立噪声-精度关系的理论分析框架",
        "为CIM芯片噪声鲁棒训练提供解决方案"
    ])

    # ========== 第18页：总结与展望 ==========
    add_text_slide(prs, "总结与展望", """总结：
• 成功实现基于STE的噪声感知训练框架
• Adaptive-STE-Sqrt性能超越基准（85.30% vs 85.15%）
• 噪声鲁棒性显著提升

未来工作：
• 在更大数据集上验证（ImageNet）
• 探索Tiki-Taka等先进模拟训练算法
• 与真实CIM硬件协同设计""")

    # ========== 第19页：致谢 ==========
    add_text_slide(prs, "感谢聆听", """感谢评委老师的指导

感谢主办方提供的平台

欢迎提问与交流""")

    return prs

if __name__ == "__main__":
    output_path = "/mnt/storage2/zyc/CIM比赛/赛道5/赛题2/文档/赛题二答辩PPT.pptx"

    print("正在生成PPT（包含实际图表）...")
    prs = create_ppt()
    prs.save(output_path)
    print(f"PPT已保存至: {output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")