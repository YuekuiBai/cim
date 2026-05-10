#!/usr/bin/env python3
"""
生成赛题二答辩PPT
基于ppt提示词.md的内容制作
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def add_title_slide(prs, title_text, subtitle_text=None):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title_text, bullet_points):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(24)
        p.space_after = Pt(16)

    return slide

def add_text_slide(prs, title_text, content_text):
    """添加文本内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(22)
        p.space_after = Pt(8)

    return slide

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== 第1页：封面 ==========
    add_title_slide(prs,
        "基于直通估计器的噪声感知训练框架设计",
        "赛道五：存算算法 - 赛题二\n2026 InnoCIM高校挑战赛")

    # ========== 第2页：目录 ==========
    add_content_slide(prs, "汇报提纲", [
        "研究背景与意义",
        "核心算法设计",
        "实验结果与分析",
        "创新点总结",
        "总结与展望"
    ])

    # ========== 第3页：研究背景 ==========
    add_content_slide(prs, "研究背景与意义", [
        "存算一体(CIM)芯片面临噪声挑战",
        "编程误差、非线性噪声、输出误差影响推理精度",
        "传统梯度下降法难以处理不可微噪声操作",
        "需要设计噪声感知的训练框架"
    ])

    # ========== 第4页：问题分析 ==========
    add_content_slide(prs, "核心问题与挑战", [
        "噪声操作的不可微性导致梯度无法直接计算",
        "传统方法无法准确估计反向传播梯度",
        "需要在\"梯度估计精度\"与\"训练可行性\"间取得平衡",
        "如何设计有效的噪声感知训练策略？"
    ])

    # ========== 第5页：STE核心思想 ==========
    add_content_slide(prs, "直通估计器(STE)核心思想", [
        "前向传播：注入真实噪声，模拟CIM硬件特性",
        "反向传播：使用STE绕过不可微操作",
        "关键洞察：梯度方向保持比精确值更重要",
        "参考CoRR 2025论文理论框架"
    ])

    # ========== 第6页：框架设计 ==========
    add_content_slide(prs, "STE噪声感知训练框架", [
        "NoiseInjector：注入编程误差、串扰、饱和非线性",
        "NoisyLinear/NoisyConv2d：带噪声的神经网络层",
        "STEGradientEstimator：自适应梯度估计",
        "支持多种噪声调度策略"
    ])

    # ========== 第7页：创新算法（一） ==========
    add_content_slide(prs, "创新算法 - 自适应STE", [
        "问题：固定缩放因子无法适应不同噪声水平",
        "解决：设计自适应梯度缩放策略",
        "四种调度：Inverse / Linear / Sqrt / Exp",
        "Sqrt调度在实验中表现最佳"
    ])

    # ========== 第8页：创新算法（二） ==========
    add_content_slide(prs, "创新算法 - 噪声感知正则化", [
        "思想：约束权重分布紧密度，降低噪声敏感性",
        "方法：L2正则化 + 噪声方差估计",
        "偏差校正：EMA方法估计并补偿梯度偏差",
        "层次化注入：根据层深度自适应调整噪声强度"
    ])

    # ========== 第9页：实验设置 ==========
    add_text_slide(prs, "实验配置", """• 数据集：CIFAR-10图像分类
• 网络架构：ResNet18
• 训练配置：30 epochs, SGD, Cosine LR
• 噪声强度：0.0 / 0.5 / 1.0 / 1.5
• 对比方法：Baseline / STE-NAT / Adaptive-STE等""")

    # ========== 第10页：核心实验结果 ==========
    add_text_slide(prs, "核心实验结果", """关键数据：
• Baseline: 85.15%
• Adaptive-STE-Sqrt: 85.30% (最佳)
• STE+Layerwise: 84.86%
• STE+BiasCorrection: 84.34%
• STE-NAT: 84.16%

自适应STE-Sqrt超越基准模型0.15%""")

    # ========== 第11页：噪声鲁棒性 ==========
    add_content_slide(prs, "噪声鲁棒性分析", [
        "噪声训练模型在干净环境下保持高性能",
        "Adaptive-STE系列方法噪声鲁棒性优于标准STE-NAT",
        "Sqrt调度策略在噪声环境下表现最稳定",
        "无噪声基准在噪声环境下性能下降约1.5%"
    ])

    # ========== 第12页：统计显著性分析 ==========
    add_text_slide(prs, "统计分析结果", """统计检验结果：
• t检验：p < 0.05，方法间存在显著差异
• ANOVA：F检验验证组间差异显著性
• 效应量Cohen's d：大于0.5为中等效应
• 95%置信区间分析各方法性能稳定性""")

    # ========== 第13页：创新点总结 ==========
    add_content_slide(prs, "主要创新点", [
        "自适应STE梯度估计：根据噪声水平动态调整缩放因子",
        "多调度策略验证：Sqrt调度效果最佳",
        "偏差校正机制：EMA方法补偿梯度估计偏差",
        "层次化噪声注入：根据层特性差异化处理"
    ])

    # ========== 第14页：技术贡献 ==========
    add_content_slide(prs, "技术贡献", [
        "提出完整的STE噪声感知训练框架",
        "验证自适应机制的有效性",
        "建立噪声-精度关系的理论分析框架",
        "为CIM芯片噪声鲁棒训练提供解决方案"
    ])

    # ========== 第15页：总结与展望 ==========
    add_text_slide(prs, "总结与展望", """总结：
• 成功实现基于STE的噪声感知训练框架
• Adaptive-STE-Sqrt性能超越基准
• 噪声鲁棒性显著提升

未来工作：
• 在更大数据集上验证（ImageNet）
• 探索Tiki-Taka等先进模拟训练算法
• 与真实CIM硬件协同设计""")

    # ========== 第16页：致谢 ==========
    add_text_slide(prs, "感谢聆听", """感谢评委老师的指导

感谢主办方提供的平台

欢迎提问与交流""")

    return prs

if __name__ == "__main__":
    output_path = "/mnt/storage2/zyc/CIM比赛/赛道5/赛题2/文档/赛题二答辩PPT.pptx"

    print("正在生成PPT...")
    prs = create_ppt()
    prs.save(output_path)
    print(f"PPT已保存至: {output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")